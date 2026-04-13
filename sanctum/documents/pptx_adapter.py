"""PowerPoint (.pptx) structured reader/writer backed by python-pptx.

Segment granularity is the *run* inside a shape's text frame, same
trade-off as the .docx adapter — formatting survives, cross-run
entities don't.

Segment IDs:
    slide{i}/shape{j}/p{p}/r{r}                       regular text-frame run
    slide{i}/shape{j}/row{r}/col{c}/p{p}/r{r}         run inside a table cell

Phase 1 scope: speaker notes, slide masters, and grouped shapes are
*not* anonymized. They flow through the raw handle untouched — lifting
that is a Phase 2 task.
"""
from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any

from pptx import Presentation

from sanctum.documents.structured import build_document, build_segment

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PptxPresentation
    from pptx.text.text import TextFrame, _Run

    from sanctum.core.models import StructuredDocument, TextSegment


def _iter_text_frame_runs(
    text_frame: TextFrame, prefix: str
) -> list[tuple[str, _Run]]:
    pairs: list[tuple[str, _Run]] = []
    for p_idx, para in enumerate(text_frame.paragraphs):
        for r_idx, run in enumerate(para.runs):
            pairs.append((f"{prefix}/p{p_idx}/r{r_idx}", run))
    return pairs


def _collect_shape_runs(shape: Any, shape_prefix: str) -> list[tuple[str, _Run]]:
    """Yield ``(segment_id, run)`` for a single shape, flattening tables."""
    pairs: list[tuple[str, _Run]] = []
    if shape.has_text_frame:
        pairs.extend(_iter_text_frame_runs(shape.text_frame, shape_prefix))
    if shape.has_table:
        for r_idx, row in enumerate(shape.table.rows):
            for c_idx, cell in enumerate(row.cells):
                prefix = f"{shape_prefix}/row{r_idx}/col{c_idx}"
                pairs.extend(_iter_text_frame_runs(cell.text_frame, prefix))
    return pairs


class Reader:
    """Flatten every slide's run-level text into segments."""

    def read(self, path: Path) -> StructuredDocument:
        prs: PptxPresentation = Presentation(str(path))
        segments: list[TextSegment] = []

        for s_idx, slide in enumerate(prs.slides):
            for sh_idx, shape in enumerate(slide.shapes):
                prefix = f"slide{s_idx}/shape{sh_idx}"
                for seg_id, run in _collect_shape_runs(shape, prefix):
                    segments.append(build_segment(seg_id, run.text))

        return build_document(
            source_path=path,
            fmt="pptx",
            segments=segments,
            raw_handle=prs,
        )


class Writer:
    """Re-apply segment text to the raw Presentation and save."""

    def write(self, doc: StructuredDocument, path: Path) -> None:
        prs: PptxPresentation | None = doc.raw_handle
        if prs is None:
            raise ValueError(
                "PptxWriter requires the StructuredDocument.raw_handle from PptxReader"
            )

        run_index = self._build_run_index(prs)
        for segment in doc.segments:
            run = run_index.get(segment.id)
            if run is None:
                continue
            run.text = segment.text

        prs.save(str(path))

    @staticmethod
    def _build_run_index(prs: PptxPresentation) -> dict[str, _Run]:
        index: dict[str, _Run] = {}
        for s_idx, slide in enumerate(prs.slides):
            for sh_idx, shape in enumerate(slide.shapes):
                prefix = f"slide{s_idx}/shape{sh_idx}"
                for seg_id, run in _collect_shape_runs(shape, prefix):
                    index[seg_id] = run
        return index
