#!/usr/bin/env python3
"""Generate deterministic office-format test fixtures for Sanctum.

Builds tiny canonical .docx / .xlsx / .pdf / .pptx files that embed the
same PII content as the existing annotated plain-text fixtures under
``tests/fixtures/{legal,financial,correspondence,...}``. The output lives
under ``tests/fixtures/office/`` and is committed to git so CI can run
offline against stable binaries.

Usage:
    python scripts/generate_office_fixtures.py

Determinism:
    Office binaries (docx/xlsx/pptx) are ZIP archives that normally embed
    creation/modification timestamps in ``docProps/core.xml`` and in each
    ZIP entry header. We post-process every generated file to normalize
    those to a fixed epoch so re-running the script produces byte-identical
    output. PDFs use reportlab's ``invariant=1`` flag for the same reason.

This script is intentionally framework-free: it does not import anything
from ``sanctum.*``. Fixtures must remain readable even if the adapter
layer is broken.
"""

from __future__ import annotations

import hashlib
import io
import json
import logging
import pathlib
import re
import sys
import zipfile
from typing import Any

from docx import Document as DocxDocument
from docx.shared import Pt
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

logger = logging.getLogger("generate_office_fixtures")

ROOT = pathlib.Path(__file__).resolve().parent.parent
FIXTURES = ROOT / "tests" / "fixtures"
OUT = FIXTURES / "office"

# Fixed epoch used for zip-entry timestamps and docProps so output is
# byte-stable across runs. Picked as 2020-01-01T00:00:00Z arbitrarily.
EPOCH = (2020, 1, 1, 0, 0, 0)
EPOCH_ISO = "2020-01-01T00:00:00Z"


# ---------------------------------------------------------------------------
# Determinism helpers
# ---------------------------------------------------------------------------


def _normalize_office_zip(path: pathlib.Path) -> None:
    """Rewrite a docx/xlsx/pptx file with fixed timestamps.

    Office formats are ZIP archives. Two sources of non-determinism:
    1. Each ZIP entry carries a ``date_time`` header.
    2. ``docProps/core.xml`` embeds ``dcterms:created`` and
       ``dcterms:modified`` elements.

    We read the archive, rewrite both, and save in place. The content is
    otherwise untouched.
    """
    buf = io.BytesIO()
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as dst:
        for info in sorted(src.infolist(), key=lambda i: i.filename):
            data = src.read(info.filename)
            if info.filename == "docProps/core.xml":
                data = _normalize_core_xml(data)
            new_info = zipfile.ZipInfo(filename=info.filename, date_time=EPOCH)
            new_info.compress_type = zipfile.ZIP_DEFLATED
            new_info.external_attr = info.external_attr
            dst.writestr(new_info, data)
    path.write_bytes(buf.getvalue())


_TIMESTAMP_RE = re.compile(
    r"(<dcterms:(created|modified)\b[^>]*>)[^<]*(</dcterms:(?:created|modified)>)"
)


def _normalize_core_xml(data: bytes) -> bytes:
    """Replace timestamps inside docProps/core.xml with the fixed epoch.

    python-docx and python-pptx emit the time elements with an outer
    ``xmlns:dcterms``/``xmlns:xsi`` declared at the root, while openpyxl
    declares both namespaces inline on each element. A regex that ignores
    the attribute payload handles all three cases.
    """
    text = data.decode("utf-8")
    text = _TIMESTAMP_RE.sub(lambda m: m.group(1) + EPOCH_ISO + m.group(3), text)
    return text.encode("utf-8")


def _sha256(path: pathlib.Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


# ---------------------------------------------------------------------------
# Source-content loader
# ---------------------------------------------------------------------------


def _load_source(rel: str) -> tuple[str, dict[str, Any]]:
    """Load a plain-text fixture and its annotation JSON."""
    txt = (FIXTURES / rel).read_text(encoding="utf-8")
    ann = json.loads((FIXTURES / rel).with_suffix(".json").read_text("utf-8"))
    return txt, ann


def _pii_inventory(annotation: dict[str, Any]) -> list[dict[str, str]]:
    """Extract a de-duplicated (entity_type, text) inventory."""
    seen: set[tuple[str, str]] = set()
    out: list[dict[str, str]] = []
    for ent in annotation["entities"]:
        key = (ent["entity_type"], ent["text"])
        if key in seen:
            continue
        seen.add(key)
        out.append({"entity_type": ent["entity_type"], "text": ent["text"]})
    return out


def _write_provenance(
    out_path: pathlib.Path,
    source_rel: str,
    format_: str,
    annotation: dict[str, Any],
) -> None:
    """Write the companion .json describing the generated fixture.

    This intentionally does NOT carry segment IDs — those depend on adapter
    implementation decisions that land later in WS2. It carries the flat PII
    inventory so an adapter test can assert "every known entity appears
    somewhere in the extracted segments" without coupling to segment naming.
    """
    provenance = {
        "document": out_path.name,
        "format": format_,
        "provenance": "synthetic",
        "generator": "scripts/generate_office_fixtures.py",
        "source_text": f"tests/fixtures/{source_rel}",
        "pii_inventory": _pii_inventory(annotation),
    }
    out_path.with_suffix(".json").write_text(
        json.dumps(provenance, indent=2) + "\n", encoding="utf-8"
    )


# ---------------------------------------------------------------------------
# Canonical fixture builders
# ---------------------------------------------------------------------------


def build_docx_nda() -> pathlib.Path:
    """NDA contract as a multi-paragraph docx with a header and a signature block."""
    source_rel = "legal/nda_contract.txt"
    text, annotation = _load_source(source_rel)
    out_path = OUT / "nda_contract.docx"

    doc = DocxDocument()

    # Header exercises the header/footer extraction path.
    section = doc.sections[0]
    section.header.paragraphs[0].text = "CONFIDENTIAL — Phase 1 fixture"

    # Body: one paragraph per line of source so extraction is non-trivial
    # and segment ids will be predictable.
    for line in text.split("\n"):
        p = doc.add_paragraph()
        if line.strip():
            run = p.add_run(line)
            run.font.size = Pt(11)

    doc.save(out_path)
    _normalize_office_zip(out_path)
    _write_provenance(out_path, source_rel, "docx", annotation)
    return out_path


def build_xlsx_invoice() -> pathlib.Path:
    """Invoice as a two-sheet workbook: 'Header' with client info, 'LineItems' with table."""
    source_rel = "financial/invoice.txt"
    _, annotation = _load_source(source_rel)
    out_path = OUT / "invoice.xlsx"

    wb = Workbook()

    # Sheet 1: client info (exercises string cells + merged header cell)
    ws1 = wb.active
    ws1.title = "Header"
    ws1["A1"] = "INVOICE"
    ws1.merge_cells("A1:D1")
    ws1["A3"] = "From:"
    ws1["B3"] = "Jenkins-Shields"
    ws1["A4"] = "Attn:"
    ws1["B4"] = "Matthew Lucas"
    ws1["A5"] = "Address:"
    ws1["B5"] = "33387 Robert Harbors Suite 317, North Deniseside, MP 02435"
    ws1["A6"] = "Phone:"
    ws1["B6"] = "432-867-7360"
    ws1["A7"] = "Email:"
    ws1["B7"] = "gallowayjoseph@example.com"
    ws1["A9"] = "To:"
    ws1["B9"] = "Henderson-Bernard"
    ws1["A10"] = "Attn:"
    ws1["B10"] = "Kristi Higgins MD"
    ws1["A11"] = "Address:"
    ws1["B11"] = "805 Brendan Neck, North Susan, CO 24857"
    ws1["A13"] = "Invoice Number:"
    ws1["B13"] = "INV-18981"
    ws1["A14"] = "Invoice Date:"
    ws1["B14"] = "March 28, 2026"
    ws1["A15"] = "Due Date:"
    ws1["B15"] = "May 15, 2026"

    # Sheet 2: line items (exercises numeric cells + one formula)
    ws2 = wb.create_sheet("LineItems")
    ws2["A1"] = "Description"
    ws2["B1"] = "Qty"
    ws2["C1"] = "Unit Price"
    ws2["D1"] = "Total"
    ws2["A2"] = "Consulting Services (40 hrs)"
    ws2["B2"] = 40
    ws2["C2"] = 150.0
    ws2["D2"] = 6000.0
    ws2["A3"] = "Project Management"
    ws2["B3"] = 1
    ws2["C3"] = 2500.0
    ws2["D3"] = 2500.0
    ws2["A4"] = "Travel Expenses"
    ws2["B4"] = 1
    ws2["C4"] = 875.50
    ws2["D4"] = 875.50
    ws2["A5"] = "TOTAL"
    ws2["D5"] = "=SUM(D2:D4)"  # formula cell: adapters must skip by default

    wb.save(out_path)
    _normalize_office_zip(out_path)
    _write_provenance(out_path, source_rel, "xlsx", annotation)
    return out_path


def build_pdf_engagement_letter() -> pathlib.Path:
    """Engagement letter as a two-page text PDF produced by reportlab."""
    source_rel = "legal/engagement_letter.txt"
    text, annotation = _load_source(source_rel)
    out_path = OUT / "engagement_letter.pdf"

    # ``invariant=1`` tells reportlab to emit byte-identical output across
    # runs (same object ids, no random xref). It also zeroes timestamps.
    doc = SimpleDocTemplate(
        str(out_path),
        pagesize=LETTER,
        title="Engagement Letter — Phase 1 fixture",
        author="sanctum-fixtures",
        creator="sanctum-fixtures",
        producer="sanctum-fixtures",
        invariant=1,
    )
    styles = getSampleStyleSheet()
    flow: list[Any] = []
    for line in text.split("\n"):
        if not line.strip():
            flow.append(Spacer(1, 8))
            continue
        # Escape the few reportlab mini-HTML chars that could trip the parser.
        safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        flow.append(Paragraph(safe, styles["BodyText"]))
    doc.build(flow)

    _write_provenance(out_path, source_rel, "pdf", annotation)
    return out_path


def build_pptx_memo() -> pathlib.Path:
    """Internal memo as a three-slide deck with speaker notes."""
    source_rel = "correspondence/internal_memo.txt"
    text, annotation = _load_source(source_rel)
    out_path = OUT / "internal_memo.pptx"

    prs = Presentation()
    # Slide 1: title
    layout_title = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(layout_title)
    slide1.shapes.title.text = "White-Green — Internal Memorandum"
    slide1.placeholders[1].text = (
        "Q2 Strategic Planning and Budget Allocation — " "From Jennifer Martin, CEO"
    )

    # Slide 2: content (bullet list) — exercises paragraph/run walking
    layout_content = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(layout_content)
    slide2.shapes.title.text = "Priorities"
    body = slide2.placeholders[1].text_frame
    body.text = "Cameron Dean — Operations, supply chain cost cut 15%"
    for para in [
        "James Elliott — Product Development, accelerate v3.0 release",
        "Marvin Taylor — CFO, revised financial forecast",
        "Follow-up meeting April 18, 2026 — Conference Room A",
    ]:
        body.add_paragraph().text = para

    # Slide 3: speaker notes (exercises notes extraction path)
    slide3 = prs.slides.add_slide(layout_content)
    slide3.shapes.title.text = "Contact"
    slide3.placeholders[1].text = "Reach Jennifer Martin at ext. 4764"
    slide3.notes_slide.notes_text_frame.text = (
        "Full memo content for reviewer reference:\n\n" + text
    )

    prs.save(out_path)
    _normalize_office_zip(out_path)
    _write_provenance(out_path, source_rel, "pptx", annotation)
    return out_path


# ---------------------------------------------------------------------------
# Entrypoint
# ---------------------------------------------------------------------------


BUILDERS = [
    build_docx_nda,
    build_xlsx_invoice,
    build_pdf_engagement_letter,
    build_pptx_memo,
]


def main() -> int:
    logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")
    OUT.mkdir(parents=True, exist_ok=True)

    manifest: dict[str, str] = {}
    for build in BUILDERS:
        path = build()
        digest = _sha256(path)
        rel = path.relative_to(ROOT)
        manifest[str(rel)] = digest
        logger.info("wrote %s (%d bytes, sha256=%s)", rel, path.stat().st_size, digest[:16])

    (OUT / "MANIFEST.sha256.json").write_text(
        json.dumps(manifest, indent=2, sort_keys=True) + "\n", encoding="utf-8"
    )
    logger.info("manifest written: %s", OUT / "MANIFEST.sha256.json")
    return 0


if __name__ == "__main__":
    sys.exit(main())
