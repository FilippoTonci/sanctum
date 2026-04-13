"""Unit tests for the .pptx reader/writer adapter pair."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from sanctum.documents.pptx_adapter import Reader, Writer

FIXTURE = Path("tests/fixtures/office/internal_memo.pptx")


@pytest.fixture()
def reader() -> Reader:
    return Reader()


@pytest.fixture()
def writer() -> Writer:
    return Writer()


def test_reader_produces_run_level_segments(reader: Reader) -> None:
    doc = reader.read(FIXTURE)
    assert doc.format == "pptx"
    assert doc.segments, "expected runs in the memo fixture"
    assert all(seg.id.startswith("slide") for seg in doc.segments)
    # Fixture contains the string 'Jennifer Martin' in slide 0 title+subtitle.
    assert any("Jennifer Martin" in seg.text for seg in doc.segments)


def test_round_trip_preserves_visible_text(reader: Reader, writer: Writer, tmp_path: Path) -> None:
    doc = reader.read(FIXTURE)
    out = tmp_path / "roundtrip.pptx"
    writer.write(doc, out)

    original = Presentation(str(FIXTURE))
    rewritten = Presentation(str(out))
    assert len(original.slides) == len(rewritten.slides)
    for a_slide, b_slide in zip(original.slides, rewritten.slides, strict=True):
        a_text = [s.text_frame.text for s in a_slide.shapes if s.has_text_frame]
        b_text = [s.text_frame.text for s in b_slide.shapes if s.has_text_frame]
        assert a_text == b_text


def test_writer_applies_segment_edits(reader: Reader, writer: Writer, tmp_path: Path) -> None:
    doc = reader.read(FIXTURE)
    target = next(seg for seg in doc.segments if "Jennifer Martin" in seg.text)
    new_segments = [
        seg.model_copy(update={"text": "<PERSON>"}) if seg.id == target.id else seg
        for seg in doc.segments
    ]
    mutated = doc.model_copy(update={"segments": new_segments})
    mutated.raw_handle = doc.raw_handle

    out = tmp_path / "edited.pptx"
    writer.write(mutated, out)

    reread = reader.read(out)
    assert any(seg.id == target.id and seg.text == "<PERSON>" for seg in reread.segments)


def test_writer_rejects_missing_raw_handle(writer: Writer, reader: Reader, tmp_path: Path) -> None:
    doc = reader.read(FIXTURE)
    doc.raw_handle = None
    with pytest.raises(ValueError, match="raw_handle"):
        writer.write(doc, tmp_path / "x.pptx")


def test_registry_dispatches_to_pptx_adapter() -> None:
    from sanctum.documents import adapter_for

    r, w = adapter_for(FIXTURE)
    assert isinstance(r, Reader)
    assert isinstance(w, Writer)
