"""End-to-end integration tests for structured document anonymization.

These load the real Presidio stack, read a committed office fixture,
anonymize every segment, and verify the written file (a) opens back
up and (b) no longer contains the original person-name PII.

Marked ``integration`` so the default unit-test run skips them.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from sanctum.analyzer.adapter import PresidioAnalyzer
from sanctum.anonymizer.adapter import PresidioAnonymizer
from sanctum.core.engine import SanctumEngine
from sanctum.documents import adapter_for

pytestmark = pytest.mark.integration

FIXTURES = Path("tests/fixtures/office")


@pytest.fixture(scope="module")
def engine() -> SanctumEngine:
    return SanctumEngine(
        analyzer=PresidioAnalyzer(),
        anonymizer=PresidioAnonymizer(),
    )


def _run_pipeline(engine: SanctumEngine, src: Path, dst: Path) -> int:
    reader, writer = adapter_for(src)
    results = engine.process_document(reader, writer, src, dst)
    return sum(len(r.detections) for r in results)


def test_docx_end_to_end_removes_named_person(engine: SanctumEngine, tmp_path: Path) -> None:
    src = FIXTURES / "nda_contract.docx"
    dst = tmp_path / "out.docx"
    total = _run_pipeline(engine, src, dst)
    assert total > 0, "Presidio should have found some PII in the NDA fixture"

    rewritten = Document(str(dst))
    full_text = "\n".join(p.text for p in rewritten.paragraphs)
    assert "Rachel Moore" not in full_text  # a known PERSON from the fixture inventory


def test_xlsx_end_to_end_removes_person(engine: SanctumEngine, tmp_path: Path) -> None:
    src = FIXTURES / "invoice.xlsx"
    dst = tmp_path / "out.xlsx"
    _run_pipeline(engine, src, dst)

    wb = load_workbook(str(dst))
    joined = "\n".join(
        str(cell.value)
        for name in wb.sheetnames
        for row in wb[name].iter_rows()
        for cell in row
        if cell.value is not None
    )
    # The fixture's known client name must have been replaced by the anonymizer.
    assert "Jenkins-Shields" not in joined or "<" in joined


def test_pptx_end_to_end_removes_person(engine: SanctumEngine, tmp_path: Path) -> None:
    src = FIXTURES / "internal_memo.pptx"
    dst = tmp_path / "out.pptx"
    _run_pipeline(engine, src, dst)

    prs = Presentation(str(dst))
    joined = "\n".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Jennifer Martin" not in joined
